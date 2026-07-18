import os
import fitz  # PyMuPDF
import logging
import tempfile
from typing import List, Dict, Any, Optional

logger = logging.getLogger(__name__)


def _ocr_scanned_page(page: "fitz.Page", page_num: int) -> Optional[Dict[str, Any]]:
    """C-3: run a scanned/image page through the multi-engine OCR orchestrator.

    The page is rendered to a temp PNG at 2x zoom (OCR accuracy needs more
    pixels than the PDF's native 72 dpi) and routed with hint="handwritten"
    so PaddleOCR is the primary engine and Docling the fallback. The
    orchestrator API is async but ingestion runs in a sync Celery task, so
    a fresh event loop is used (same pattern as proactive insights).

    Returns {"text": ..., "engine": ..., "confidence": ...} or None on
    failure — the caller decides the fallback and MUST log it loudly.
    """
    import asyncio
    from app.services.ocr_orchestrator import ocr_orchestrator

    tmp_path = None
    try:
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            tmp_path = tmp.name
        pix.save(tmp_path)

        loop = asyncio.new_event_loop()
        try:
            result = loop.run_until_complete(
                ocr_orchestrator.extract_document(
                    tmp_path, mime_type="image/png", hint="handwritten"
                )
            )
        finally:
            loop.close()

        if not result.text.strip() or result.text.strip() == "Mock Text":
            return None
        return {
            "text": result.text,
            "engine": result.engine_name,
            "confidence": result.confidence,
        }
    except Exception as exc:
        logger.error(
            f"[OCR] Orchestrator failed on scanned page {page_num}: {exc}"
        )
        return None
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.remove(tmp_path)


# P8 — PowerPoint extraction. python-pptx is an optional dep; gracefully
# degrade if not installed so PDF-only deployments still work.
try:
    from pptx import Presentation  # type: ignore
    _PPTX_AVAILABLE = True
except Exception:  # pragma: no cover — env without python-pptx
    Presentation = None  # type: ignore
    _PPTX_AVAILABLE = False


def _extract_pptx_stream(file_path: str):
    """Yield one page-equivalent record per slide.

    Slide text is concatenated from every shape that exposes a text frame
    (title placeholder, body, text boxes, tables). Slide notes are appended
    after a `--- speaker notes ---` divider so the chunker can decide
    whether to include them in retrieval.
    """
    if not _PPTX_AVAILABLE or Presentation is None:
        raise RuntimeError(
            "python-pptx is not installed; .pptx extraction unavailable."
        )
    logger.info(f"[Tracing] Starting python-pptx extraction for {file_path}")
    prs = Presentation(file_path)
    for idx, slide in enumerate(prs.slides):
        parts: List[str] = []
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            parts.append(line)
                # tables in slides
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        joined = " | ".join(c for c in cells if c)
                        if joined:
                            parts.append(joined)
            except Exception as inner_exc:
                logger.debug(
                    f"[pptx] shape extraction skipped on slide {idx + 1}: {inner_exc}"
                )
        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes_text = ""

        slide_text = "\n\n".join(parts).strip()
        if notes_text:
            slide_text = (
                f"{slide_text}\n\n--- speaker notes ---\n{notes_text}".strip()
            )

        yield {
            "page_number": idx + 1,  # treat each slide as a "page" for citations
            "extracted_text": slide_text,
            "layout_metadata": {
                "source": "pptx",
                "slide_index": idx,
                "shape_count": len(slide.shapes),
                "has_notes": bool(notes_text),
            },
        }
    logger.info(f"[Tracing] python-pptx extraction complete: {idx + 1} slides")


class OCRService:
    @staticmethod
    def is_text_native(page: fitz.Page) -> bool:
        """
        Heuristic to detect if a page is text-native vs scanned.
        Returns False if there is very little text (likely an image-only scan).
        """
        text = page.get_text("text")
        return len(text.strip()) > 50

    @staticmethod
    def extract_layout_metadata(page: fitz.Page) -> Dict[str, Any]:
        """
        Extracts block-level metadata, crucial for semantic chunking later.
        """
        blocks = page.get_text("dict").get("blocks", [])
        text_blocks = [b for b in blocks if b.get("type") == 0]
        
        return {
            "blocks_count": len(text_blocks),
            "width": page.rect.width,
            "height": page.rect.height
        }

    @staticmethod
    def _looks_like_pptx(file_path: str) -> bool:
        """P8: detect a pptx via extension OR ZIP-magic + ppt/ directory.

        document_tasks.py downloads every storage object to a
        `tempfile(...suffix=".pdf")` regardless of the original mime, so
        the extension alone isn't reliable — peek at the bytes too.
        PPTX files are ZIP archives whose central directory contains
        `ppt/presentation.xml`.
        """
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pptx":
            return True
        try:
            with open(file_path, "rb") as fh:
                head = fh.read(4)
            if head[:2] != b"PK":
                return False
            import zipfile
            with zipfile.ZipFile(file_path) as zf:
                names = zf.namelist()
            return any(n.startswith("ppt/") for n in names)
        except Exception:
            return False

    @staticmethod
    def extract_document_stream(file_path: str):
        """
        Extracts document pages using layout-aware mechanisms, yielding page by page to save memory.

        P8: dispatches by detection. `.pptx` files are streamed through
        python-pptx (one yield per slide); everything else goes through
        the existing PyMuPDF path so PDFs keep their behaviour.
        """
        if OCRService._looks_like_pptx(file_path):
            yield from _extract_pptx_stream(file_path)
            return

        logger.info(f"[Tracing] Starting PyMuPDF streaming extraction for {file_path}")

        try:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                is_native = OCRService.is_text_native(page)
                
                if not is_native:
                    # C-3: scanned/image pages go through the real OCR
                    # orchestrator (PaddleOCR primary, Docling fallback).
                    # OCR_SCANNED_ENABLED is the rollback toggle: heavy
                    # engines can be disabled without a code change.
                    from app.core.config import settings
                    ocr_result = None
                    if settings.OCR_SCANNED_ENABLED:
                        logger.info(
                            f"[Tracing] Page {page_num+1} appears scanned — routing to OCR orchestrator."
                        )
                        ocr_result = _ocr_scanned_page(page, page_num + 1)
                    if ocr_result:
                        extracted_text = ocr_result["text"]
                        layout_meta = {
                            "is_native": False,
                            "ocr_engine": ocr_result["engine"],
                            "ocr_confidence": ocr_result["confidence"],
                        }
                    else:
                        # Degraded path — loud, observable, never silent.
                        logger.error(
                            f"[Tracing] Page {page_num+1} is scanned and OCR was "
                            f"{'disabled' if not settings.OCR_SCANNED_ENABLED else 'unavailable/failed'} — "
                            "falling back to raw text extraction (likely empty)."
                        )
                        extracted_text = page.get_text("text")
                        layout_meta = {
                            "is_native": False,
                            "requires_fallback": True,
                            "ocr_failed": settings.OCR_SCANNED_ENABLED,
                        }
                else:
                    # Layout-aware extraction preserving reading order
                    blocks = page.get_text("blocks")
                    # PyMuPDF blocks format: (x0, y0, x1, y1, "text", block_no, block_type)
                    # Sort roughly top-to-bottom, left-to-right
                    sorted_blocks = sorted(blocks, key=lambda b: (b[1], b[0]))
                    
                    # Extract only text blocks (type == 0)
                    clean_text = "\n\n".join([b[4].strip() for b in sorted_blocks if b[6] == 0])
                    
                    layout_meta = OCRService.extract_layout_metadata(page)
                    layout_meta["is_native"] = True
                    extracted_text = clean_text

                yield {
                    "page_number": page_num + 1,  # 1-indexed for citation systems
                    "extracted_text": extracted_text,
                    "layout_metadata": layout_meta
                }
                
            doc.close()
            logger.info(f"[Tracing] Completed streaming extraction.")
            
        except Exception as e:
            logger.error(f"[Tracing] Extraction failed for {file_path}: {str(e)}")
            raise e
